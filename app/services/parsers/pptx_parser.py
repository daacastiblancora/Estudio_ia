from typing import List
from langchain_core.documents import Document
from app.services.parser_registry import BaseParser
from app.core.logging import logger

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class PPTXParser(BaseParser):
    """PowerPoint parser — one Document per slide with title tracking."""
    supported_extensions = ["pptx"]

    def parse(self, file_path: str, source_name: str) -> List[Document]:
        if Presentation is None:
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            return []

        docs = []
        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides, 1):
            texts = []
            title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)

                    # Capture title shape
                    if hasattr(shape, "name") and "Title" in shape.name:
                        title = shape.text_frame.text.strip()

            if texts:
                content = self._inject_header(source_name, "Diapositiva", idx)
                if title:
                    content += f"Título: {title}\n\n"
                content += "\n".join(texts)

                docs.append(Document(
                    page_content=content,
                    metadata=self._make_metadata(
                        source=source_name, page=idx, total=total_slides,
                        fmt="pptx", section=title
                    ),
                ))

        return docs
